import re
import csv
from pathlib import Path

import pdfplumber
from docx import Document

# functions to extract text from each file type
def extract_pdf(path: str):
    text = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text

def extract_docx(path: str):
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def extract_txt(path: str):
    return Path(path).read_text(encoding="utf-8", errors="ignore")

# markdown is just plain text for extraction purposes — the ##/**/etc.
# syntax doesn't hurt question generation, no need to strip it
def extract_md(path: str):
    return Path(path).read_text(encoding="utf-8", errors="ignore")

def extract_rtf(path: str):
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        raise ValueError(
            "RTF support requires the 'striprtf' package. "
            "Install it with: pip install striprtf"
        )
    raw = Path(path).read_text(encoding="utf-8", errors="ignore")
    return rtf_to_text(raw)

def extract_pptx(path: str):
    try:
        from pptx import Presentation
    except ImportError:
        raise ValueError(
            "PPTX support requires the 'python-pptx' package. "
            "Install it with: pip install python-pptx"
        )
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)

def extract_csv(path: str):
    rows = []
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            if any(cell.strip() for cell in row):
                rows.append(", ".join(row))
    return "\n".join(rows)

def extract_html(path: str):
    raw = Path(path).read_text(encoding="utf-8", errors="ignore")
    # strip script/style blocks entirely first so their contents don't
    # leak into the extracted text, then strip remaining tags
    raw = re.sub(r"<(script|style)[^>]*>.*?</\1>", " ", raw, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", raw)
    text = re.sub(r"&nbsp;", " ", text)
    text = re.sub(r"&amp;", "&", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()

def extract_text(path: str):
    ext = Path(path).suffix.lower().lstrip(".")
    extractors = {
        "pdf": extract_pdf,
        "docx": extract_docx,
        "txt": extract_txt,
        "md": extract_md,
        "rtf": extract_rtf,
        "pptx": extract_pptx,
        "csv": extract_csv,
        "html": extract_html,
        "htm": extract_html,
    }
    if ext not in extractors:
        raise ValueError(f"Unsupported file type: .{ext}")
    return extractors[ext](path)

# breaks text into chunks to be individually processed in order to generate questions on more specific sections
def chunk_text(text: str, max_chars: int):
    if len(text) <= max_chars:
        return [text]

    sentences = re.split(r"(?<=[.!?])\s+", text)
    chunks = []
    current = ""
    for sentence in sentences:
        if len(current) + len(sentence) + 1 > max_chars:
            if current:
                chunks.append(current.strip())
            current = sentence
        else:
            current = (current + " " + sentence).strip()
    if current:
        chunks.append(current.strip())
    return chunks

# filters for acknowledgement pages and publisher info
acknowledgement_filters = [
    r"all rights reserved",
    r"copyright\s*©?\s*\d{4}",
    r"no part of this (publication|book|document) may be reproduced",
    r"isbn[\s:]*[\d-]+",
    r"printed in",
    r"prior written permission",
    r"brief quotations used in reviews",

    r"published by",
    r"publisher:",
    r"first edition",
    r"www\.\S+\.(com|org|net)",
    r"\d{3,}\s+\w+\s+(avenue|street|road|blvd|drive)",  # street addresses
    r"email:\s*\S+@\S+",
    r"website:\s*www\.",

    r"acknowledg(e)?ments?",
    r"would like to express their sincere",
    r"special thanks",
    r"we also acknowledge",
    r"we thank our families",
    r"throughout the preparation of this",
    r"whose feedback has been invaluable",
    r"inspired the creation of this",

    r"this document is intended to demonstrate",
    r"content (contained in this document|is entirely fictional)",
    r"provided solely for (educational|testing|demonstration)",
    r"any resemblance to actual",
    r"we hope this sample document",
    r"serves as a useful reference",
    r"table of contents",
    r"library of congress",
]

def remove_acknowledgement(text):
    sentences = re.split(r"(?<=[.!?])\s+", text)
    filtered = []

    for s in sentences:
        s_lower = s.lower()
        if any(re.search(p, s_lower) for p in acknowledgement_filters):
            continue
        filtered.append(s)

    return " ".join(filtered)